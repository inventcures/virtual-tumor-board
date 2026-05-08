"""
Virtual Tumor Board — clinician-facing deck for CMC Vellore.

Design principles applied (Saloni Dattani's data-viz guide):
  - Headline-style titles convey the takeaway, not just the topic
  - Plain language; jargon defined inline
  - Direct labels (no legends), horizontal text only, no 3D, no chart junk
  - Logical category ordering (workflow phases left-to-right; specialties grouped)
  - Colour used semantically: red = caveat, green = benefit, blue = information,
    slate = body text
  - Every content slide carries a source/footer line; deck is standalone
  - Single sans-serif (Calibri); generous whitespace; no decorative clip-art
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Palette (semantic only) -------------------------------------------------
INK       = RGBColor(0x1A, 0x23, 0x32)   # body text — near-black slate
MUTED     = RGBColor(0x6B, 0x77, 0x85)   # captions, footers
RULE      = RGBColor(0xCB, 0xD2, 0xD9)   # hairlines
PAGE_BG   = RGBColor(0xFF, 0xFF, 0xFF)
INFO      = RGBColor(0x2B, 0x6C, 0xB0)   # blue — information / accent
BENEFIT   = RGBColor(0x2F, 0x85, 0x5D)   # green — safety / benefit
CAVEAT    = RGBColor(0xC5, 0x30, 0x30)   # red — risk / caveat / "what it is NOT"
PHASE_BG  = RGBColor(0xEE, 0xF2, 0xF7)   # very light blue — phase boxes
PANEL_BG  = RGBColor(0xF7, 0xF8, 0xFA)   # very light slate — panels

FONT = "Calibri"

# --- Layout constants (16:9, 13.333" x 7.5") --------------------------------
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
TITLE_TOP = Inches(0.45)
TITLE_H   = Inches(0.95)
KICKER_TOP = Inches(0.30)
BODY_TOP  = Inches(1.55)
FOOTER_TOP = Inches(7.05)
FOOTER_H   = Inches(0.30)


# --- Helpers ----------------------------------------------------------------
def set_text(tf, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             margin=Pt(2), word_wrap=True):
    """Replace tf contents with a list of (text, {style}) tuples — one para each."""
    tf.clear()
    tf.word_wrap = word_wrap
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = margin
    tf.margin_bottom = margin
    tf.vertical_anchor = anchor
    for i, item in enumerate(runs):
        text, style = item
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = style.get("align", align)
        if "space_after" in style:
            para.space_after = style["space_after"]
        if "level" in style:
            para.level = style["level"]
        # Support inline runs: text can be a list of (sub_text, sub_style).
        if isinstance(text, list):
            for j, (sub_text, sub_style) in enumerate(text):
                run = para.runs[0] if (j == 0 and para.runs) else para.add_run()
                run.text = sub_text
                _apply_run(run, {**style, **sub_style})
        else:
            run = para.add_run()
            run.text = text
            _apply_run(run, style)


def _apply_run(run, style):
    f = run.font
    f.name = style.get("font", FONT)
    f.size = style.get("size", Pt(14))
    f.bold = style.get("bold", False)
    f.italic = style.get("italic", False)
    f.color.rgb = style.get("color", INK)


def add_textbox(slide, left, top, width, height, runs, **kwargs):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box.text_frame, runs, **kwargs)
    return box


def add_rect(slide, left, top, width, height, fill, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    shp.text_frame.text = ""
    return shp


def add_hairline(slide, left, top, width, color=RULE):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(0.75)
    return line


def add_title(slide, kicker, title):
    """Kicker = small uppercase section tag. Title = full-sentence takeaway."""
    add_textbox(
        slide, MARGIN_L, KICKER_TOP, CONTENT_W, Inches(0.30),
        [(kicker.upper(), {"size": Pt(11), "bold": True, "color": INFO})],
    )
    add_textbox(
        slide, MARGIN_L, TITLE_TOP + Inches(0.20), CONTENT_W, TITLE_H,
        [(title, {"size": Pt(26), "bold": True, "color": INK})],
    )
    add_hairline(slide, MARGIN_L, BODY_TOP - Inches(0.15), CONTENT_W)


def add_footer(slide, source_text):
    add_hairline(slide, MARGIN_L, FOOTER_TOP - Inches(0.05), CONTENT_W)
    add_textbox(
        slide, MARGIN_L, FOOTER_TOP, CONTENT_W, FOOTER_H,
        [(source_text, {"size": Pt(9), "color": MUTED, "italic": True})],
    )


def bullets(items, *, size=Pt(15), space=Pt(8), color=INK):
    """Convert a list of strings (or (label,desc) tuples) into runs for set_text."""
    runs = []
    for item in items:
        if isinstance(item, tuple):
            label, desc = item
            runs.append((
                [(f"{label}  ", {"bold": True, "color": INK}),
                 (desc, {"color": color})],
                {"size": size, "space_after": space, "color": color},
            ))
        else:
            runs.append((f"•  {item}",
                         {"size": size, "space_after": space, "color": color}))
    return runs


# =============================================================================
# Build deck
# =============================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------------------
# Slide 1 — Title
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
# left rule
add_rect(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, INFO)

add_textbox(
    s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(0.5),
    [("VIRTUAL TUMOR BOARD",
      {"size": Pt(13), "bold": True, "color": INFO})],
)
add_textbox(
    s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.6),
    [("An AI co-pilot for multidisciplinary cancer deliberation",
      {"size": Pt(38), "bold": True, "color": INK})],
)
add_textbox(
    s, Inches(0.9), Inches(4.20), Inches(11.5), Inches(0.6),
    [("Seven AI specialists debate every case, three meta-reviewers push back, "
      "one moderator synthesises a structured recommendation.",
      {"size": Pt(17), "color": MUTED, "italic": True})],
)

add_hairline(s, Inches(0.9), Inches(5.55), Inches(8.0))
add_textbox(
    s, Inches(0.9), Inches(5.65), Inches(11.5), Inches(0.4),
    [("Presented to the Department of Oncology, CMC Vellore  ·  May 2026",
      {"size": Pt(13), "color": INK})],
)
add_textbox(
    s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.4),
    [("Live demo:  virtual-tumor-board-production.up.railway.app",
      {"size": Pt(12), "color": INFO})],
)

# ----------------------------------------------------------------------------
# Slide 2 — Origin story (the founder's aunt)
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title(s, "Why this exists",
          "An ambiguous MRI line decided whether one patient's cancer was curable.")

# Two columns: the case on the left, the wider lesson on the right.
left_w  = Inches(7.0)
right_w = CONTENT_W - left_w - Inches(0.4)
top     = BODY_TOP + Inches(0.05)
panel_h = Inches(5.0)

# Left panel — the personal trigger (story)
add_rect(s, MARGIN_L, top, left_w, panel_h, PANEL_BG)
add_rect(s, MARGIN_L, top, Inches(0.10), panel_h, INFO)
add_textbox(
    s, MARGIN_L + Inches(0.30), top + Inches(0.20), left_w - Inches(0.5), Inches(0.45),
    [("THE CASE", {"size": Pt(11), "bold": True, "color": INFO})],
)
add_textbox(
    s, MARGIN_L + Inches(0.30), top + Inches(0.70), left_w - Inches(0.5),
    panel_h - Inches(0.85),
    [
        (
            "The founder's aunt was diagnosed with cancer. Her MRI showed a "
            "vertebral lesion. The radiology report read:",
            {"size": Pt(13.5), "color": INK, "space_after": Pt(10)},
        ),
        (
            "\u201csclerotic change, possibly age-related, "
            "cannot rule out metastatic disease.\u201d",
            {"size": Pt(14), "italic": True, "color": INK,
             "space_after": Pt(10)},
        ),
        (
            "That single line decided whether her case was Stage II "
            "(curable) or Stage IV (incurable).",
            {"size": Pt(13.5), "color": INK, "space_after": Pt(14)},
        ),
        ([
            ("Through family connections, ",
             {"color": INK}),
            ("we got her case reviewed by top specialists.",
             {"color": INK, "bold": True}),
            (" They concluded it was benign. She received curative treatment.",
             {"color": INK}),
        ], {"size": Pt(13.5), "space_after": Pt(10)}),
        ([
            ("Most patients aren't that lucky.", {"color": CAVEAT, "bold": True}),
            (" They don't have connections. They can't travel to metros. "
             "They can't afford multiple specialist consultations.",
             {"color": INK}),
        ], {"size": Pt(13.5)}),
    ],
)

# Right panel — the systemic gap
right_x = MARGIN_L + left_w + Inches(0.4)
add_rect(s, right_x, top, right_w, panel_h, PHASE_BG)
add_rect(s, right_x, top, Inches(0.10), panel_h, CAVEAT)
add_textbox(
    s, right_x + Inches(0.30), top + Inches(0.20), right_w - Inches(0.5), Inches(0.45),
    [("THE GAP", {"size": Pt(11), "bold": True, "color": CAVEAT})],
)
add_textbox(
    s, right_x + Inches(0.30), top + Inches(0.70), right_w - Inches(0.5),
    panel_h - Inches(0.85),
    [
        ([("< 5%", {"size": Pt(34), "bold": True, "color": CAVEAT})],
         {"space_after": Pt(2)}),
        ("of cancer patients in India receive multidisciplinary expert "
         "review — the gold standard for optimal outcomes.",
         {"size": Pt(12.5), "color": INK, "space_after": Pt(14)}),

        ([("1.4 M", {"size": Pt(28), "bold": True, "color": INK})],
         {"space_after": Pt(2)}),
        ("new cancer cases annually.",
         {"size": Pt(12.5), "color": INK, "space_after": Pt(14)}),

        ([("4\u00d7 fewer", {"size": Pt(24), "bold": True, "color": INK})],
         {"space_after": Pt(2)}),
        ("oncologists per capita than the United States — concentrated in "
         "metros; tier-2/3 cities face delayed and suboptimal care.",
         {"size": Pt(12.5), "color": INK}),
    ],
)

add_textbox(
    s, MARGIN_L, FOOTER_TOP - Inches(0.55), CONTENT_W, Inches(0.4),
    [("VTB exists so the next patient like her doesn't need a connection to "
      "get a second opinion.",
      {"size": Pt(12.5), "italic": True, "color": INFO, "bold": True})],
)
add_footer(s, "Source: inventcures.github.io/ai-powered-virtual-mtb  ·  "
              "Founder: Ashish (tp53)  ·  Stage: working prototype, not yet approved for clinical use.")

# ----------------------------------------------------------------------------
# Slide 3 — The clinical problem
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title(s, "The problem",
          "Tumor boards work — but they don't scale to every case, every week.")

add_textbox(
    s, MARGIN_L, BODY_TOP, CONTENT_W, Inches(4.8),
    bullets([
        ("At a tertiary centre,",
         "only a fraction of new cancer cases are formally reviewed at MDT — "
         "the rest are decided in single-specialty clinics."),
        ("Boards are constrained by",
         "clinician time, room scheduling, and sub-specialty availability "
         "(genetics, palliative, radiation oncology aren't always at the table)."),
        ("Trainees rarely",
         "see the full reasoning chain across specialties — it's heard, not "
         "documented, and rarely revisited."),
        ("Yet NCCN, ESMO and the WHO position MDT review",
         "as the standard of care for almost every solid tumour."),
        ("The gap is not knowledge — it is throughput.",
         "An AI co-pilot can stage a deliberation in minutes, on every case, "
         "before the human MDT meets."),
    ], size=Pt(15), space=Pt(12)),
)
add_footer(s, "Sources: NCCN Guidelines (2025); ESMO Clinical Practice Guidelines; "
              "WHO Cancer Control Knowledge into Action.")

# ----------------------------------------------------------------------------
# Slide 3 — What VTB is, in one screen (3 columns)
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title(s, "What VTB does",
          "VTB stages a multidisciplinary deliberation in minutes, on any case.")

col_w = (CONTENT_W - Inches(0.4)) / 3
col_top = BODY_TOP
col_h = Inches(4.8)

def feature_panel(left, header, lines, accent=INFO):
    add_rect(s, left, col_top, col_w, col_h, PANEL_BG)
    add_rect(s, left, col_top, Inches(0.08), col_h, accent)  # left accent rule
    add_textbox(
        s, left + Inches(0.30), col_top + Inches(0.30), col_w - Inches(0.5),
        Inches(0.5),
        [(header, {"size": Pt(17), "bold": True, "color": INK})],
    )
    add_textbox(
        s, left + Inches(0.30), col_top + Inches(0.95), col_w - Inches(0.5),
        col_h - Inches(1.0),
        [(line, {"size": Pt(13), "color": INK, "space_after": Pt(8)})
         for line in lines],
    )

feature_panel(
    MARGIN_L,
    "Multi-agent debate",
    [
        "•  7 AI oncology specialists + 3 meta-reviewers",
        "•  5-phase Chain-of-Debate (gatekeeper → opinions → critique → rebuttal → consensus)",
        "•  Powered by Claude Sonnet 4 with retrieval over NCCN/ESMO guidelines",
    ],
    accent=INFO,
)
feature_panel(
    MARGIN_L + col_w + Inches(0.2),
    "Imaging analysis",
    [
        "•  MedGemma 27B for radiology reads (CXR, CT, MRI, mammography)",
        "•  Modality-aware prompts with explicit confidence scores",
        "•  Accepts DICOM and a phone-camera photograph of a printed film",
    ],
    accent=BENEFIT,
)
feature_panel(
    MARGIN_L + 2 * (col_w + Inches(0.2)),
    "Live, auditable output",
    [
        "•  Server-Sent-Event stream — clinicians watch each agent reason in real time",
        "•  Structured consensus with citations and recorded dissent",
        "•  Open-source (MIT) and self-hostable behind a hospital firewall",
    ],
    accent=INFO,
)
add_footer(s, "Live: virtual-tumor-board-production.up.railway.app  ·  "
              "Code: github.com/inventcures/virtual-tumor-board (MIT).")

# ----------------------------------------------------------------------------
# Slide 4 — The team in a box (specialists + meta-agents)
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title(s, "The team",
          "Seven specialists make the call. Three meta-reviewers keep them honest.")

# Two side-by-side panels
left_w  = Inches(7.6)
right_w = CONTENT_W - left_w - Inches(0.3)
panel_top = BODY_TOP
panel_h = Inches(5.0)

# Specialists (left)
add_rect(s, MARGIN_L, panel_top, left_w, panel_h, PANEL_BG)
add_textbox(
    s, MARGIN_L + Inches(0.3), panel_top + Inches(0.20), left_w - Inches(0.6), Inches(0.5),
    [("CLINICAL SPECIALISTS  ·  7 agents, run in parallel",
      {"size": Pt(11), "bold": True, "color": INFO})],
)
specialists = [
    ("Dr. Shalya",     "Surgical Oncology — resectability, surgical approach"),
    ("Dr. Chikitsa",   "Medical Oncology — systemic therapy, sequencing"),
    ("Dr. Kirann",     "Radiation Oncology — protocol, dose, fractionation"),
    ("Dr. Chitran",    "Onco-Radiology — staging, response assessment"),
    ("Dr. Marga",      "Pathology — histology, IHC, molecular markers"),
    ("Dr. Anuvamsha",  "Genetics — germline & somatic, hereditary syndromes"),
    ("Dr. Shanti",     "Palliative Care — symptom burden, goals of care"),
]
add_textbox(
    s, MARGIN_L + Inches(0.3), panel_top + Inches(0.65), left_w - Inches(0.6),
    panel_h - Inches(0.8),
    [([(f"{name}", {"bold": True, "color": INK}),
       (f"   {role}", {"color": INK})],
      {"size": Pt(13.5), "space_after": Pt(8)})
     for name, role in specialists],
)

# Meta-agents (right)
add_rect(s, MARGIN_L + left_w + Inches(0.3), panel_top, right_w, panel_h, PHASE_BG)
add_textbox(
    s, MARGIN_L + left_w + Inches(0.6), panel_top + Inches(0.20),
    right_w - Inches(0.6), Inches(0.5),
    [("META-REVIEWERS  ·  applied after each round",
      {"size": Pt(11), "bold": True, "color": CAVEAT})],
)
meta = [
    ("Dr. Tark",
     "Scientific Critic — challenges evidence, demands citations"),
    ("Dr. Samata",
     "Stewardship — equity, cost, antibiotic / chemo stewardship"),
    ("Principal Investigator",
     "Moderator — gatekeeps completeness, synthesises consensus, surfaces dissent"),
]
add_textbox(
    s, MARGIN_L + left_w + Inches(0.6), panel_top + Inches(0.65),
    right_w - Inches(0.9), panel_h - Inches(0.8),
    [([(f"{name}", {"bold": True, "color": INK}),
       (f"\n{role}", {"color": INK})],
      {"size": Pt(13.5), "space_after": Pt(14)})
     for name, role in meta],
)

add_footer(s, "Source: packages/agents/src/specialists/index.ts  ·  "
              "All agents run on Claude Sonnet 4 (claude-sonnet-4-20250514) with tool-use for guideline retrieval.")

# ----------------------------------------------------------------------------
# Slide 5 — The 5-phase deliberation flow
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title(s, "How a deliberation runs",
          "Specialists debate, critics push back, the moderator synthesises.")

phases = [
    ("Phase 0",  "Gatekeeper",
     "PI checks the case is complete enough to deliberate.", INFO),
    ("Phase 1",  "Initial opinions  (parallel)",
     "All 7 specialists respond simultaneously. Each can call retrieval tools for guidelines.", INFO),
    ("Phase 2",  "Critique  (sequential)",
     "Scientific Critic and Stewardship reviewer challenge the Round-1 plans.", CAVEAT),
    ("Phase 3",  "Round 2  —  rebuttal & resolution",
     "PI moderates a focused debate where conflicts are reconciled.", INFO),
    ("Phase 4",  "Consensus",
     "PI synthesises a structured recommendation, with citations and recorded dissent.", BENEFIT),
]

box_h = Inches(0.85)
box_gap = Inches(0.20)
box_left = MARGIN_L
box_top = BODY_TOP + Inches(0.2)
label_w = Inches(1.1)
title_w = Inches(3.5)
desc_w  = CONTENT_W - label_w - title_w - Inches(0.4)

for i, (label, ttl, desc, color) in enumerate(phases):
    y = box_top + i * (box_h + box_gap)
    # left coloured tab
    add_rect(s, box_left, y, Inches(0.18), box_h, color)
    # background panel
    add_rect(s, box_left + Inches(0.18), y, CONTENT_W - Inches(0.18), box_h, PANEL_BG)
    # phase label
    add_textbox(
        s, box_left + Inches(0.40), y + Inches(0.18), label_w, Inches(0.5),
        [(label, {"size": Pt(12), "bold": True, "color": color})],
    )
    # phase title
    add_textbox(
        s, box_left + Inches(0.40) + label_w, y + Inches(0.13), title_w, Inches(0.6),
        [(ttl, {"size": Pt(15), "bold": True, "color": INK})],
    )
    # phase description
    add_textbox(
        s, box_left + Inches(0.40) + label_w + title_w, y + Inches(0.18),
        desc_w, Inches(0.6),
        [(desc, {"size": Pt(12.5), "color": INK})],
    )

# Annotation underneath
add_textbox(
    s, MARGIN_L, FOOTER_TOP - Inches(0.55), CONTENT_W, Inches(0.4),
    [("Live SSE stream throughout — clinicians watch the reasoning unfold, "
      "not just the verdict.",
      {"size": Pt(12), "italic": True, "color": INFO})],
)
add_footer(s, "Source: packages/agents/src/orchestrator/index.ts (V7 spec, lines 68–152); "
              "apps/web/src/app/api/deliberate/stream/route.ts (SSE).")

# ----------------------------------------------------------------------------
# Slide 6 — MedGemma imaging
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title(s, "Imaging — radiology reads",
          "Modality-aware reads with explicit confidence and target-lesion measurements.")

# Pipeline boxes
steps = [
    ("Image in",
     "DICOM, NIfTI, or a JPG/PNG photo of a printed film"),
    ("Modality-specific prompt",
     "Different prompt template for CXR, CT, MRI, mammography"),
    ("MedGemma 27B",
     "Google's medical foundation model, called via Vertex AI"),
    ("Structured output",
     "Findings (location, severity), RECIST-style measurements (≥10 mm targets), "
     "confidence score (0.30–0.95), model attribution"),
]
step_w = (CONTENT_W - Inches(0.6)) / 4
step_top = BODY_TOP + Inches(0.2)
step_h = Inches(2.4)

for i, (ttl, body) in enumerate(steps):
    x = MARGIN_L + i * (step_w + Inches(0.2))
    add_rect(s, x, step_top, step_w, step_h, PANEL_BG)
    add_rect(s, x, step_top, step_w, Inches(0.08), INFO)
    add_textbox(
        s, x + Inches(0.20), step_top + Inches(0.25), step_w - Inches(0.4), Inches(0.5),
        [(ttl, {"size": Pt(14), "bold": True, "color": INK})],
    )
    add_textbox(
        s, x + Inches(0.20), step_top + Inches(0.85), step_w - Inches(0.4),
        step_h - Inches(1.0),
        [(body, {"size": Pt(12), "color": INK})],
    )
    if i < 3:
        # arrow between boxes (simple chevron)
        arrow = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x + step_w + Inches(0.02),
            step_top + step_h / 2 - Inches(0.1),
            Inches(0.16), Inches(0.20),
        )
        arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED
        arrow.line.fill.background()

# Reliability note
add_rect(s, MARGIN_L, step_top + step_h + Inches(0.3),
         CONTENT_W, Inches(0.55), PHASE_BG)
add_textbox(
    s, MARGIN_L + Inches(0.25), step_top + step_h + Inches(0.36),
    CONTENT_W - Inches(0.5), Inches(0.5),
    [([("Three-tier fallback so a deliberation never stalls:  ",
        {"bold": True, "color": INK}),
       ("Vertex AI MedGemma  →  HuggingFace MedGemma Space  →  Gemini 2.0 Flash.",
        {"color": INK})],
      {"size": Pt(12.5)})],
)

# Caveat
add_rect(s, MARGIN_L, step_top + step_h + Inches(1.0),
         CONTENT_W, Inches(0.55), RGBColor(0xFD, 0xEC, 0xEC))
add_textbox(
    s, MARGIN_L + Inches(0.25), step_top + step_h + Inches(1.07),
    CONTENT_W - Inches(0.5), Inches(0.5),
    [([("Caveat.  ", {"bold": True, "color": CAVEAT}),
       ("MedGemma is a research model, not a regulated diagnostic device. "
        "Use as decision support; the radiologist remains the source of truth.",
        {"color": INK})],
      {"size": Pt(12.5)})],
)
add_footer(s, "Source: apps/web/src/lib/medgemma/client.ts (lines 49–143, 564–650).  "
              "RECIST = Response Evaluation Criteria In Solid Tumours.")

# ----------------------------------------------------------------------------
# Slide 8 — Inputs that work in a real Indian tertiary centre
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title(s, "Inputs",
          "Works with what's already in radiology workflow — and with a phone.")

inputs_data = [
    ("DICOM series",
     ".dcm files or full study folders",
     "CT, MRI, X-ray, PET, ultrasound, bone scan",
     "Drag-and-drop the entire study; parsed in the browser using dicom-parser."),
    ("Phone camera",
     "JPEG / PNG via device camera",
     "When PACS is not at hand",
     "Photograph a printed film at the bedside, on a ward round, or in an outreach clinic."),
    ("Gallery upload",
     "JPEG / PNG from device storage",
     "Prior scans the patient brought on a phone",
     "Useful when families bring outside imaging that never made it onto PACS."),
]

row_h = Inches(1.05)
top = BODY_TOP + Inches(0.1)
header_y = top
add_rect(s, MARGIN_L, header_y, CONTENT_W, Inches(0.45), PHASE_BG)
col_xs = [MARGIN_L + Inches(0.20),
          MARGIN_L + Inches(2.6),
          MARGIN_L + Inches(5.2),
          MARGIN_L + Inches(8.5)]
col_ws = [Inches(2.4), Inches(2.6), Inches(3.3), Inches(3.5)]
headers = ["Input", "Format", "When to use it", "What VTB does with it"]
for x, w, h in zip(col_xs, col_ws, headers):
    add_textbox(
        s, x, header_y + Inches(0.08), w, Inches(0.35),
        [(h.upper(), {"size": Pt(10), "bold": True, "color": INFO})],
    )

for i, row in enumerate(inputs_data):
    y = header_y + Inches(0.45) + i * row_h
    if i % 2 == 1:
        add_rect(s, MARGIN_L, y, CONTENT_W, row_h, PANEL_BG)
    for x, w, val in zip(col_xs, col_ws, row):
        add_textbox(
            s, x, y + Inches(0.18), w, row_h - Inches(0.2),
            [(val, {"size": Pt(12.5), "color": INK})],
        )

# Annotation
add_textbox(
    s, MARGIN_L, FOOTER_TOP - Inches(0.55), CONTENT_W, Inches(0.4),
    [("Critical for ward rounds, peripheral OPDs, and outreach clinics where "
      "PACS access is limited.",
      {"size": Pt(12), "italic": True, "color": INFO})],
)
add_footer(s, "Source: apps/web/src/components/my-imaging/DicomUploader.tsx; "
              "ImagingUploadPanel.tsx.  Imaging path does not accept PDFs.")

# ----------------------------------------------------------------------------
# Slide 9 — Architecture & deployment
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title(s, "Architecture",
          "A single Next.js app on Railway, talking to specialised model APIs.")

# Three layers stacked
layer_left = MARGIN_L
layer_w = CONTENT_W
layer_top = BODY_TOP + Inches(0.05)
layer_h = Inches(1.35)
gap = Inches(0.25)

def layer(idx, kicker, items, accent):
    y = layer_top + idx * (layer_h + gap)
    add_rect(s, layer_left, y, layer_w, layer_h, PANEL_BG)
    add_rect(s, layer_left, y, Inches(0.10), layer_h, accent)
    add_textbox(
        s, layer_left + Inches(0.30), y + Inches(0.12), Inches(2.6), Inches(0.4),
        [(kicker.upper(), {"size": Pt(10), "bold": True, "color": accent})],
    )
    # items as side-by-side cards inside the layer
    n = len(items)
    inner_left = layer_left + Inches(2.95)
    inner_w = layer_w - Inches(3.1)
    card_w = (inner_w - Inches(0.2) * (n - 1)) / n
    for j, (head, sub) in enumerate(items):
        cx = inner_left + j * (card_w + Inches(0.2))
        add_textbox(
            s, cx, y + Inches(0.18), card_w, Inches(0.4),
            [(head, {"size": Pt(13.5), "bold": True, "color": INK})],
        )
        add_textbox(
            s, cx, y + Inches(0.62), card_w, layer_h - Inches(0.7),
            [(sub, {"size": Pt(11.5), "color": INK})],
        )

layer(0, "Client",
      [("Next.js 15 + React 19", "Browser app — uploads, deliberation UI, "
                                  "DICOM parsing all client-side"),
       ("Auth (optional)",       "SITE_ACCESS_TOKEN cookie gate — single shared "
                                  "token; no patient login layer yet")],
      INFO)

layer(1, "Application — on Railway (NIXPACKS)",
      [("Next.js API routes",     "REST endpoints, Server Actions, SSE stream "
                                  "for live deliberation"),
       ("Agent orchestrator",     "5-phase Chain-of-Debate over 10 agents; "
                                  "tool-use loop for guideline retrieval"),
       ("PostgreSQL (Railway)",   "Analytics only — visitor IP, geolocation, "
                                  "feature usage. No patient data persisted.")],
      BENEFIT)

layer(2, "External model APIs",
      [("Anthropic Claude Sonnet 4", "All 10 reasoning agents"),
       ("Vertex AI MedGemma 27B",    "Radiology reads (HF + Gemini fallback)"),
       ("Gemini File Search",        "RAG retrieval over NCCN, ESMO, ASTRO, "
                                      "ACR, CAP, SSO, ClinVar/CIViC")],
      INFO)

add_footer(s, "Source: railway.json; apps/web/src/middleware.ts.  "
              "Stack: Next.js 15, React 19, TailwindCSS, pnpm workspaces.")

# ----------------------------------------------------------------------------
# Slide 10 — Three CMC Vellore workflows
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title(s, "How CMC Vellore could use VTB",
          "Three concrete entry points — none require new infrastructure.")

workflows = [
    ("A.  Pre-board preparation",
     "Junior registrar runs VTB on each case the night before the weekly MDT.",
     "Brings the structured consensus, the dissenting opinions, and the "
     "retrieved guideline citations to the in-person board. The human MDT "
     "spends its time on disagreement, not recap.",
     BENEFIT),
    ("B.  Outreach clinic / ward round",
     "Senior takes a phone photograph of a printed scan a referring centre "
     "sent with the patient.",
     "Within minutes, the system returns a radiology read, a tentative stage, "
     "and a treatment-plan sketch — enough to triage the patient that visit "
     "instead of asking them to return.",
     INFO),
    ("C.  Trainee education",
     "Residents watch the live SSE stream of agents debating an archived case "
     "between weekly boards.",
     "Builds MDT intuition: how does a medical oncologist push back against a "
     "surgical oncologist? Where does the stewardship reviewer object? Visible "
     "reasoning, not just a verdict.",
     INFO),
]

card_w = (CONTENT_W - Inches(0.4)) / 3
top = BODY_TOP + Inches(0.05)
card_h = Inches(4.7)
for i, (head, who, what, accent) in enumerate(workflows):
    x = MARGIN_L + i * (card_w + Inches(0.2))
    add_rect(s, x, top, card_w, card_h, PANEL_BG)
    add_rect(s, x, top, card_w, Inches(0.10), accent)
    add_textbox(
        s, x + Inches(0.25), top + Inches(0.30), card_w - Inches(0.5), Inches(0.6),
        [(head, {"size": Pt(15.5), "bold": True, "color": INK})],
    )
    add_textbox(
        s, x + Inches(0.25), top + Inches(1.05), card_w - Inches(0.5), Inches(1.2),
        [([("Who & when.  ", {"bold": True, "color": MUTED}),
           (who, {"color": INK})], {"size": Pt(12.5)})],
    )
    add_textbox(
        s, x + Inches(0.25), top + Inches(2.55), card_w - Inches(0.5),
        card_h - Inches(2.7),
        [([("What VTB adds.  ", {"bold": True, "color": MUTED}),
           (what, {"color": INK})], {"size": Pt(12.5)})],
    )

add_textbox(
    s, MARGIN_L, FOOTER_TOP - Inches(0.55), CONTENT_W, Inches(0.4),
    [("In every workflow, VTB augments the human MDT — it never replaces it.",
      {"size": Pt(12.5), "italic": True, "color": BENEFIT, "bold": True})],
)
add_footer(s, "These workflows are illustrative — each would need a small pilot at CMC Vellore before any wider use.")

# ----------------------------------------------------------------------------
# Slide 11 — Safety, governance, what VTB is NOT
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title(s, "Safety & governance",
          "A decision-support tool — not a diagnostic device, not a regulatory product.")

# Two columns: what it is / what it is not
left_w  = (CONTENT_W - Inches(0.4)) / 2
right_w = left_w
top = BODY_TOP + Inches(0.05)
panel_h = Inches(5.0)

# Left — what it IS
add_rect(s, MARGIN_L, top, left_w, panel_h, PANEL_BG)
add_rect(s, MARGIN_L, top, Inches(0.10), panel_h, BENEFIT)
add_textbox(
    s, MARGIN_L + Inches(0.30), top + Inches(0.20), left_w - Inches(0.5), Inches(0.45),
    [("WHAT VTB IS", {"size": Pt(11), "bold": True, "color": BENEFIT})],
)
is_items = [
    ("Decision support.",
     "Every recommendation comes with citations, dissenting agent opinions, "
     "and an explicit confidence on imaging reads."),
    ("Stateless on patient data.",
     "Case text and images are sent to the model APIs at request time. "
     "Nothing about an individual patient is stored in the application database."),
    ("Auditable.",
     "The full SSE transcript can be saved as a record of how a recommendation "
     "was reached."),
    ("Open-source (MIT).",
     "The institution can read the code, change it, or self-host it behind the "
     "hospital firewall."),
]
add_textbox(
    s, MARGIN_L + Inches(0.30), top + Inches(0.75), left_w - Inches(0.5),
    panel_h - Inches(0.85),
    [([(label + "  ", {"bold": True, "color": INK}),
       (body, {"color": INK})],
      {"size": Pt(13), "space_after": Pt(12)})
     for label, body in is_items],
)

# Right — what it is NOT
right_x = MARGIN_L + left_w + Inches(0.4)
add_rect(s, right_x, top, right_w, panel_h, RGBColor(0xFD, 0xEC, 0xEC))
add_rect(s, right_x, top, Inches(0.10), panel_h, CAVEAT)
add_textbox(
    s, right_x + Inches(0.30), top + Inches(0.20), right_w - Inches(0.5), Inches(0.45),
    [("WHAT VTB IS NOT", {"size": Pt(11), "bold": True, "color": CAVEAT})],
)
isnot_items = [
    ("Not a diagnostic device.",
     "No CDSCO, FDA, or CE clearance. Research and education use only."),
    ("Not a substitute for the radiologist or pathologist.",
     "MedGemma reads and OncoSeg masks must be reviewed before any clinical use."),
    ("Not a substitute for the human MDT.",
     "It accelerates and structures preparation. The clinical decision still "
     "happens in the room."),
    ("Not yet hardened for PHI.",
     "Today's hosted instance runs on Railway with US-region model APIs. "
     "Any institutional rollout would need a self-hosted, privacy-reviewed deployment."),
]
add_textbox(
    s, right_x + Inches(0.30), top + Inches(0.75), right_w - Inches(0.5),
    panel_h - Inches(0.85),
    [([(label + "  ", {"bold": True, "color": INK}),
       (body, {"color": INK})],
      {"size": Pt(13), "space_after": Pt(12)})
     for label, body in isnot_items],
)

add_footer(s, "PHI = Protected Health Information.  CDSCO = Central Drugs Standard Control Organisation (India).")

# ----------------------------------------------------------------------------
# Slide 12 — Limitations, roadmap, how to try it
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title(s, "Limitations, roadmap, and next step",
          "Honest about what's missing — and a small pilot is the obvious next step.")

# Three columns
col_w = (CONTENT_W - Inches(0.4)) / 3
top = BODY_TOP + Inches(0.05)
col_h = Inches(4.5)

def column(left, kicker, items, accent):
    add_rect(s, left, top, col_w, col_h, PANEL_BG)
    add_rect(s, left, top, col_w, Inches(0.08), accent)
    add_textbox(
        s, left + Inches(0.25), top + Inches(0.20), col_w - Inches(0.5), Inches(0.4),
        [(kicker.upper(), {"size": Pt(11), "bold": True, "color": accent})],
    )
    add_textbox(
        s, left + Inches(0.25), top + Inches(0.65), col_w - Inches(0.5),
        col_h - Inches(0.8),
        [(f"•  {it}", {"size": Pt(13), "color": INK, "space_after": Pt(10)})
         for it in items],
    )

column(
    MARGIN_L,
    "Limitations today",
    [
        "English-only deliberations",
        "30–90 s latency per deliberation",
        "No EHR or PACS integration yet",
        "No structured outcome capture for audit",
        "Not yet approved for clinical use",
    ],
    CAVEAT,
)
column(
    MARGIN_L + col_w + Inches(0.2),
    "Roadmap",
    [
        "PACS / DICOM-Web ingest",
        "Self-hosted, on-prem deployment for hospitals",
        "Vernacular language support (Tamil, Hindi)",
        "Outcome capture, so we can audit recommendations against follow-up",
        "Site-specific guideline overlays (CMC, NCG India)",
    ],
    INFO,
)
column(
    MARGIN_L + 2 * (col_w + Inches(0.2)),
    "Try it / talk to me",
    [
        "Live demo:  virtual-tumor-board-production.up.railway.app",
        "Code (MIT):  github.com/inventcures/virtual-tumor-board",
        "Sample report:  v18 ovarian tumour-board PDF in /sample_ouputs",
        "Contact:  acgt0101@gmail.com",
        "Proposed next step:  a 10-case retrospective pilot at CMC, blinded against the in-person MDT decision.",
    ],
    BENEFIT,
)

add_footer(s, "Slides follow Saloni Dattani's data-visualisation guide "
              "(scientificdiscovery.dev/p/salonis-guide-to-data-visualization).")

# ----------------------------------------------------------------------------
out = "Virtual_Tumor_Board__CMC_Vellore.pptx"
prs.save(out)
print(f"Wrote {out}")
